"""
SVG → PPTX Converter (Phase 2)
Converts SVG files/slides to editable PowerPoint shapes.
No rasterization — all elements become native PPTX objects.

Phase 2 enhancements:
- Gradient fill (linearGradient, radialGradient)
- Opacity / alpha
- Better text box sizing
- Corner radius adjustment
- <defs> / <use> support
- Stroke-dasharray
- Line cap / join
- <image> support (base64 embedded)
"""

import sys
import os
import re
import base64
import io
import xml.etree.ElementTree as ET
from typing import Optional, Union
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

sys.path.insert(0, os.path.dirname(__file__))
from svg_path_parser import parse_path_d

# ─── Constants ────────────────────────────────────────────
SVG_NS = 'http://www.w3.org/2000/svg'
PX_TO_EMU = 9525
SLIDE_W = int(1280 * PX_TO_EMU)
SLIDE_H = int(720 * PX_TO_EMU)

NAMED_COLORS = {
    'black': (0,0,0), 'white': (255,255,255), 'red': (255,0,0),
    'green': (0,128,0), 'blue': (0,0,255), 'gray': (128,128,128),
    'grey': (128,128,128), 'yellow': (255,255,0), 'cyan': (0,255,255),
    'magenta': (255,0,255), 'orange': (255,165,0), 'purple': (128,0,128),
    'pink': (255,192,203), 'brown': (165,42,42), 'navy': (0,0,128),
    'teal': (0,128,128), 'maroon': (128,0,0), 'lime': (0,255,0),
    'olive': (128,128,0), 'silver': (192,192,192), 'gold': (255,215,0),
    'coral': (255,127,80), 'salmon': (250,128,114),
    'transparent': None,
}


# ─── Color utilities ──────────────────────────────────────

def parse_color(color_str: str) -> Optional[RGBColor]:
    if not color_str or color_str in ('none', 'transparent', ''):
        return None
    m = re.match(r'^#([0-9a-fA-F]{6})$', color_str)
    if m:
        return RGBColor(int(m.group(1)[0:2], 16), int(m.group(1)[2:4], 16), int(m.group(1)[4:6], 16))
    m = re.match(r'^#([0-9a-fA-F]{3})$', color_str)
    if m:
        return RGBColor(int(m.group(1)[0], 16)*17, int(m.group(1)[1], 16)*17, int(m.group(1)[2], 16)*17)
    m = re.match(r'^rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)$', color_str)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    c = NAMED_COLORS.get(color_str.lower().strip())
    return RGBColor(*c) if c else None


def hex_to_rgb_tuple(hex_str: str) -> Optional[tuple[int, int, int]]:
    """Parse color string to (r,g,b) tuple."""
    c = parse_color(hex_str)
    return (c[0], c[1], c[2]) if c else None


def lerp_color(c1: tuple, c2: tuple, t: float) -> tuple:
    return tuple(int(a + (b - a) * t) for a, b in zip(c1, c2))


# ─── SVG utility ──────────────────────────────────────────

def px(val, default: float = 0) -> int:
    if val is None:
        return int(default * PX_TO_EMU)
    return int(float(val) * PX_TO_EMU)


def parse_transform(transform_str: str) -> dict:
    result = {'translate_x': 0, 'translate_y': 0, 'rotate': 0,
              'scale_x': 1, 'scale_y': 1, 'rotate_cx': None, 'rotate_cy': None}
    if not transform_str:
        return result
    for m in re.finditer(r'(translate|rotate|scale|matrix)\s*\(([^)]+)\)', transform_str):
        cmd, args_s = m.group(1), m.group(2)
        args = [float(x) for x in re.findall(r'[-+]?(?:\d+\.?\d*|\.\d+)', args_s)]
        if cmd == 'translate':
            result['translate_x'] += args[0] * PX_TO_EMU
            if len(args) > 1:
                result['translate_y'] += args[1] * PX_TO_EMU
        elif cmd == 'rotate':
            result['rotate'] += args[0]
            if len(args) >= 3:
                result['rotate_cx'] = args[1]
                result['rotate_cy'] = args[2]
        elif cmd == 'scale':
            result['scale_x'] *= args[0]
            result['scale_y'] *= args[1] if len(args) > 1 else args[0]
        elif cmd == 'matrix' and len(args) == 6:
            a, b, c, d, e, f = args
            result['translate_x'] += e * PX_TO_EMU
            result['translate_y'] += f * PX_TO_EMU
            import math
            result['rotate'] += math.degrees(math.atan2(b, a))
            result['scale_x'] *= math.sqrt(a*a + b*b)
            result['scale_y'] *= math.sqrt(c*c + d*d)
    return result


def get_attr(el, name: str, default: str = '') -> str:
    """Get attribute from element, checking both attribute and style."""
    val = el.get(name, '')
    if not val:
        style = el.get('style', '')
        m = re.search(rf'{name}:\s*([^;]+)', style)
        if m:
            val = m.group(1).strip().strip("'\"")
    return val or default


# ─── Gradient support ─────────────────────────────────────

class GradientDef:
    """Parsed SVG gradient definition."""
    def __init__(self, gid: str, gtype: str, colors: list[tuple[float, str]],
                 x1=0, y1=0, x2=1, y2=1, cx=0.5, cy=0.5, r=0.5):
        self.id = gid
        self.type = gtype  # 'linear' or 'radial'
        self.colors = colors  # [(offset, color_hex), ...]
        self.x1, self.y1, self.x2, self.y2 = x1, y1, x2, y2
        self.cx, self.cy, self.r = cx, cy, r


def _pct(v):
    v = v.strip()
    if v.endswith('%'):
        return float(v[:-1]) / 100.0
    return float(v)

def parse_gradients(root) -> dict[str, GradientDef]:
    """Extract all gradient definitions from SVG <defs>."""
    gradients = {}
    ns = SVG_NS
    
    for defs in root.findall(f'{{{ns}}}defs'):
        for lg in defs.findall(f'{{{ns}}}linearGradient'):
            gid = lg.get('id', '')
            stops = []
            for stop in lg.findall(f'{{{ns}}}stop'):
                offset = float(stop.get('offset', '0').rstrip('%')) / 100
                color = get_attr(stop, 'stop-color', '#000000')
                stops.append((offset, color))
            gradients[gid] = GradientDef(
                gid, 'linear', stops,
                x1=_pct(lg.get('x1', '0')), y1=_pct(lg.get('y1', '0')),
                x2=_pct(lg.get('x2', '1')), y2=_pct(lg.get('y2', '0'))
            )
        
        for rg in defs.findall(f'{{{ns}}}radialGradient'):
            gid = rg.get('id', '')
            stops = []
            for stop in rg.findall(f'{{{ns}}}stop'):
                offset = float(stop.get('offset', '0').rstrip('%')) / 100
                color = get_attr(stop, 'stop-color', '#000000')
                stops.append((offset, color))
            gradients[gid] = GradientDef(
                gid, 'radial', stops,
                cx=_pct(rg.get('cx', '0.5')), cy=_pct(rg.get('cy', '0.5')),
                r=_pct(rg.get('r', '0.5'))
            )
    
    return gradients


def apply_gradient_fill(shape, gradient: GradientDef):
    """Apply gradient fill to PPTX shape using XML manipulation."""
    from pptx.oxml import parse_xml
    from lxml import etree
    
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = shape._element.makeelement(qn('p:spPr'), {})
        shape._element.append(spPr)
    
    # Remove existing fill
    for old_fill in list(spPr):
        tag = old_fill.tag.split('}')[-1] if '}' in old_fill.tag else old_fill.tag
        if tag in ('solidFill', 'gradFill', 'noFill'):
            spPr.remove(old_fill)
    
    if gradient.type == 'linear':
        # Build OOXML linear gradient
        # angle in OOXML is in 60000ths of a degree, 0 = left-to-right
        import math
        dx = gradient.x2 - gradient.x1
        dy = gradient.y2 - gradient.y1
        angle_deg = math.degrees(math.atan2(-dy, dx))  # SVG: y-down, PPTX: y-up
        angle_60k = int(angle_deg * 60000)
        
        grad_fill = etree.SubElement(spPr, qn('a:gradFill'))
        gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
        
        for offset, color_hex in gradient.colors:
            gs = etree.SubElement(gs_lst, qn('a:gs'))
            gs.set('pos', str(int(offset * 100000)))
            srgb = etree.SubElement(gs, qn('a:srgbClr'))
            c = parse_color(color_hex)
            srgb.set('val', f'{c[0]:02X}{c[1]:02X}{c[2]:02X}' if c else '000000')
        
        lin = etree.SubElement(grad_fill, qn('a:lin'))
        lin.set('ang', str(angle_60k))
        lin.set('scaled', '1')
    
    elif gradient.type == 'radial':
        grad_fill = etree.SubElement(spPr, qn('a:gradFill'))
        gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
        
        for offset, color_hex in gradient.colors:
            gs = etree.SubElement(gs_lst, qn('a:gs'))
            gs.set('pos', str(int(offset * 100000)))
            srgb = etree.SubElement(gs, qn('a:srgbClr'))
            c = parse_color(color_hex)
            srgb.set('val', f'{c[0]:02X}{c[1]:02X}{c[2]:02X}' if c else '000000')
        
        path = etree.SubElement(grad_fill, qn('a:path'))
        path.set('path', 'circle')
        fill_to = etree.SubElement(path, qn('a:fillToRect'))
        fill_to.set('l', str(int(gradient.cx * 100000 - gradient.r * 50000)))
        fill_to.set('t', str(int(gradient.cy * 100000 - gradient.r * 50000)))
        fill_to.set('r', str(int(gradient.cx * 100000 + gradient.r * 50000)))
        fill_to.set('b', str(int(gradient.cy * 100000 + gradient.r * 50000)))


# ─── Opacity support ──────────────────────────────────────

def apply_opacity(shape, opacity_str: Optional[str]):
    """Apply opacity to shape fill via XML alpha channel."""
    if not opacity_str:
        return
    try:
        opacity = float(opacity_str)
    except (ValueError, TypeError):
        return
    
    if opacity >= 1.0:
        return
    
    alpha = int(opacity * 100000)
    
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    
    solid_fill = spPr.find(qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is not None:
            # Add alpha element
            from lxml import etree
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(alpha))


# ─── Corner radius ────────────────────────────────────────

def set_corner_radius(shape, rx_px: float, shape_w_emu: int, shape_h_emu: int):
    """Set rounded rectangle corner radius via XML."""
    # PPTX uses adjLst with adj value = radius / min(width, height) * 50000
    min_dim = min(shape_w_emu, shape_h_emu)
    if min_dim == 0:
        return
    radius_emu = int(rx_px * PX_TO_EMU)
    adj_val = min(int(radius_emu / min_dim * 50000), 50000)
    
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return
    
    adjLst = prstGeom.find(qn('a:adjLst'))
    if adjLst is None:
        from lxml import etree
        adjLst = etree.SubElement(prstGeom, qn('a:adjLst'))
    
    adj = adjLst.find(qn('a:adj'))
    if adj is None:
        from lxml import etree
        adj = etree.SubElement(adjLst, qn('a:adj'))
    
    adj.set('val', str(adj_val))


# ─── Image support ────────────────────────────────────────

def add_image(slide, el, sx, sy, tx, ty):
    """SVG <image> → PPTX picture shape (supports base64 data URIs)."""
    href = el.get('href', '') or el.get(f'{{{SVG_NS}}}href', '') or el.get('{http://www.w3.org/1999/xlink}href', '')
    if not href:
        return
    
    x = px(el.get('x', '0')) + tx
    y = px(el.get('y', '0')) + ty
    w = px(el.get('width', '0'))
    h = px(el.get('height', '0'))
    
    x = int(x * sx / PX_TO_EMU) if sx != 1 else x
    y = int(y * sy / PX_TO_EMU) if sy != 1 else y
    w = int(w * sx / PX_TO_EMU) if sx != 1 else w
    h = int(h * sy / PX_TO_EMU) if sy != 1 else h
    
    if w == 0 or h == 0:
        return
    
    try:
        if href.startswith('data:'):
            # base64 data URI
            header, data = href.split(',', 1)
            img_bytes = base64.b64decode(data)
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, x, y, w, h)
        else:
            # File path (relative)
            if os.path.exists(href):
                slide.shapes.add_picture(href, x, y, w, h)
    except Exception as e:
        print(f"  ⚠ Image add failed: {e}")


# ─── Stroke helpers ───────────────────────────────────────

def apply_stroke_extras(shape, el):
    """Apply stroke-dasharray, stroke-linecap, stroke-linejoin."""
    style = el.get('style', '')
    
    # Dash pattern
    dash = el.get('stroke-dasharray', '')
    if not dash:
        m = re.search(r'stroke-dasharray:\s*([^;]+)', style)
        if m:
            dash = m.group(1).strip()
    
    if dash and dash != 'none':
        # Map common dash patterns to PPTX preset dash
        line = shape.line
        if hasattr(line, '_ln'):
            from lxml import etree
            ln = line._ln
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            if '5,5' in dash or '5 5' in dash:
                prstDash.set('val', 'dash')
            elif '2,2' in dash or '2 2' in dash:
                prstDash.set('val', 'dot')
            elif '10,5' in dash or '10 5' in dash:
                prstDash.set('val', 'dashLong')
            else:
                prstDash.set('val', 'dash')


# ─── Main converter class ─────────────────────────────────

class SvgToPptx:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank_layout = self.prs.slide_layouts[6]
        self._gradients: dict[str, GradientDef] = {}
        self._defs_elements: dict[str, ET.Element] = {}
    
    def add_svg_slide(self, svg_string: str, bg_color: Optional[str] = None):
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        if bg_color:
            bg = parse_color(bg_color)
            if bg:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = bg
        
        try:
            root = ET.fromstring(svg_string)
        except ET.ParseError as e:
            print(f"  ⚠ SVG parse error: {e}")
            return
        
        vb = root.get('viewBox', '0 0 1280 720').split()
        vb_w, vb_h = float(vb[2]), float(vb[3])
        sx = SLIDE_W / vb_w
        sy = SLIDE_H / vb_h
        
        # Parse gradients from <defs>
        self._gradients = parse_gradients(root)
        
        # Parse defs elements for <use> references
        self._parse_defs(root)
        
        # Background from SVG
        self._extract_bg(root, slide, sx, sy)
        
        # Process elements
        for el in root:
            self._process(slide, el, sx, sy, 0, 0)
    
    def _parse_defs(self, root):
        """Parse <defs> for reusable elements."""
        self._defs_elements = {}
        ns = SVG_NS
        for defs in root.findall(f'{{{ns}}}defs'):
            for el in defs:
                eid = el.get('id', '')
                if eid:
                    self._defs_elements[eid] = el
    
    def _extract_bg(self, root, slide, sx, sy):
        style = root.get('style', '')
        m = re.search(r'background(?:-color)?:\s*([^;]+)', style)
        if m:
            bg = parse_color(m.group(1).strip())
            if bg:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = bg
    
    def _resolve_use(self, el) -> Optional[ET.Element]:
        """Resolve <use> href to actual element."""
        href = el.get('href', '') or el.get(f'{{{SVG_NS}}}href', '')
        if href.startswith('#'):
            return self._defs_elements.get(href[1:])
        return None
    
    def _process(self, slide, el, sx, sy, tx, ty, parent_opacity: float = 1.0):
        tag = el.tag.replace(f'{{{SVG_NS}}}', '')
        tf = parse_transform(el.get('transform', ''))
        total_tx = tx + tf['translate_x']
        total_ty = ty + tf['translate_y']
        
        # Apply scale from transform
        local_sx = sx * tf['scale_x']
        local_sy = sy * tf['scale_y']
        
        # Inherit opacity from parent <g> or own opacity
        own_opacity = get_attr(el, 'opacity', '')
        if own_opacity:
            try:
                local_opacity = parent_opacity * float(own_opacity)
            except (ValueError, TypeError):
                local_opacity = parent_opacity
        else:
            local_opacity = parent_opacity
        
        # Store opacity on element for _apply_fill / _add_text to use
        if local_opacity < 1.0:
            if not el.get('opacity'):
                el.set('opacity', str(local_opacity))
            # Also set in style for get_attr to find
            style = el.get('style', '')
            if 'opacity' not in style:
                el.set('style', f'{style};opacity:{local_opacity}' if style else f'opacity:{local_opacity}')
        
        if tag == 'rect':
            self._add_rect(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'circle':
            self._add_circle(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'ellipse':
            self._add_ellipse(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'text':
            self._add_text(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'line':
            self._add_line(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'path':
            self._add_path(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'g':
            for child in el:
                self._process(slide, child, local_sx, local_sy, total_tx, total_ty, local_opacity)
        elif tag in ('polygon', 'polyline'):
            self._add_poly(slide, el, local_sx, local_sy, total_tx, total_ty, tag == 'polygon')
        elif tag == 'image':
            add_image(slide, el, local_sx, local_sy, total_tx, total_ty)
        elif tag == 'use':
            ref = self._resolve_use(el)
            if ref is not None:
                # Create a copy of the referenced element with position offset
                import copy
                ref_copy = copy.deepcopy(ref)
                ux = float(el.get('x', '0'))
                uy = float(el.get('y', '0'))
                self._process(slide, ref_copy, local_sx, local_sy,
                              total_tx + int(ux * local_sx), total_ty + int(uy * local_sy), local_opacity)
        # Skip: defs, style, title, desc, metadata
    
    # ─── Fill / Stroke application ─────────────────────────
    
    def _get_gradient_ref(self, fill_str: str) -> Optional[GradientDef]:
        """Extract gradient ID from url(#id) reference."""
        m = re.match(r'url\(#(.+)\)', fill_str)
        if m:
            return self._gradients.get(m.group(1))
        return None
    
    def _apply_fill(self, shape, el):
        fill_str = get_attr(el, 'fill', 'none')
        
        # Check for gradient
        gradient = self._get_gradient_ref(fill_str)
        if gradient:
            apply_gradient_fill(shape, gradient)
            return
        
        if fill_str and fill_str != 'none':
            c = parse_color(fill_str)
            if c:
                shape.fill.solid()
                shape.fill.fore_color.rgb = c
            else:
                shape.fill.background()
        else:
            shape.fill.background()
        
        # Apply opacity
        opacity = get_attr(el, 'opacity')
        if not opacity:
            style = el.get('style', '')
            m = re.search(r'opacity:\s*([^;]+)', style)
            if m:
                opacity = m.group(1).strip()
        apply_opacity(shape, opacity)
    
    def _apply_stroke(self, shape, el):
        stroke_str = get_attr(el, 'stroke', 'none')
        
        if stroke_str and stroke_str != 'none':
            c = parse_color(stroke_str)
            if c:
                shape.line.color.rgb = c
                sw = get_attr(el, 'stroke-width', '1')
                shape.line.width = Pt(float(sw))
                
                # Stroke opacity
                stroke_opacity = get_attr(el, 'stroke-opacity')
                if stroke_opacity:
                    try:
                        so = float(stroke_opacity)
                        if so < 1.0:
                            alpha = int(so * 100000)
                            from lxml import etree
                            ln = shape.line._ln
                            solidFill = ln.find(qn('a:solidFill'))
                            if solidFill is not None:
                                srgb = solidFill.find(qn('a:srgbClr'))
                                if srgb is not None:
                                    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                                    alpha_el.set('val', str(alpha))
                    except (ValueError, TypeError):
                        pass
                
                # Dash pattern etc
                apply_stroke_extras(shape, el)
        else:
            shape.line.fill.background()
    
    # ─── Shape builders ────────────────────────────────────
    
    def _add_rect(self, slide, el, sx, sy, tx, ty):
        x = px(el.get('x', '0')) + tx
        y = px(el.get('y', '0')) + ty
        w = px(el.get('width', '0'))
        h = px(el.get('height', '0'))
        rx = float(el.get('rx', '0') or '0')
        ry = float(el.get('ry', '0') or '0')
        
        if w == 0 or h == 0:
            return
        
        x = int(x * sx / PX_TO_EMU) if sx != 1 else x
        y = int(y * sy / PX_TO_EMU) if sy != 1 else y
        w = int(w * sx / PX_TO_EMU) if sx != 1 else w
        h = int(h * sy / PX_TO_EMU) if sy != 1 else h
        
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (rx > 0 or ry > 0) else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        
        # Corner radius
        if rx > 0:
            set_corner_radius(shape, rx, w, h)
        
        self._apply_fill(shape, el)
        self._apply_stroke(shape, el)
    
    def _add_circle(self, slide, el, sx, sy, tx, ty):
        cx_v = float(el.get('cx', '0'))
        cy_v = float(el.get('cy', '0'))
        r_v = float(el.get('r', '0'))
        
        x = int((cx_v - r_v) * sx) + tx
        y = int((cy_v - r_v) * sy) + ty
        w = int(2 * r_v * sx)
        h = int(2 * r_v * sy)
        
        if w == 0 or h == 0:
            return
        
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        self._apply_fill(shape, el)
        self._apply_stroke(shape, el)
    
    def _add_ellipse(self, slide, el, sx, sy, tx, ty):
        cx_v = float(el.get('cx', '0'))
        cy_v = float(el.get('cy', '0'))
        rx_v = float(el.get('rx', '0'))
        ry_v = float(el.get('ry', '0'))
        
        x = int((cx_v - rx_v) * sx) + tx
        y = int((cy_v - ry_v) * sy) + ty
        w = int(2 * rx_v * sx)
        h = int(2 * ry_v * sy)
        
        if w == 0 or h == 0:
            return
        
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        self._apply_fill(shape, el)
        self._apply_stroke(shape, el)
    
    def _cjk_count(self, text: str) -> int:
        """Count CJK characters in text."""
        return sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f'
                   or '\uff00' <= c <= '\uffef' or '\u3400' <= c <= '\u4dbf')

    def _add_rect(self, slide, el, sx, sy, tx, ty):
        x = px(el.get('x', '0')) + tx
        y = px(el.get('y', '0')) + ty
        w = px(el.get('width', '0'))
        h = px(el.get('height', '0'))
        rx = float(el.get('rx', '0') or '0')
        ry = float(el.get('ry', '0') or '0')
        
        if w == 0 or h == 0:
            return
        
        x = int(x * sx / PX_TO_EMU) if sx != 1 else x
        y = int(y * sy / PX_TO_EMU) if sy != 1 else y
        w = int(w * sx / PX_TO_EMU) if sx != 1 else w
        h = int(h * sy / PX_TO_EMU) if sy != 1 else h
        
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (rx > 0 or ry > 0) else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        
        # Corner radius
        if rx > 0:
            set_corner_radius(shape, rx, w, h)
        
        self._apply_fill(shape, el)
        self._apply_stroke(shape, el)

    def _add_text(self, slide, el, sx, sy, tx, ty):
        """Convert SVG <text> element to PPTX textbox.
        
        Key fixes vs original:
        1. text-anchor:middle → shift text box left by width/2
        2. SVG y is baseline → shift up by ~fontSize*0.8 to get top
        3. Each tspan with x/y/dy gets its own textbox (SVG absolute positioning)
        4. Width estimation from position + slide edge, not fragile _find_card_bounds
        5. Inherited opacity from parent <g> elements
        """
        x_val = float(el.get('x', '0'))
        y_val = float(el.get('y', '0'))
        
        font_size = float(get_attr(el, 'font-size', '18'))
        font_family = get_attr(el, 'font-family', 'Arial')
        font_weight = get_attr(el, 'font-weight', 'normal')
        fill_color = get_attr(el, 'fill', '#000000')
        text_anchor = get_attr(el, 'text-anchor', 'start')
        opacity = get_attr(el, 'opacity', '')
        if not opacity:
            style = el.get('style', '')
            m = re.search(r'opacity:\s*([^;]+)', style)
            if m:
                opacity = m.group(1).strip()
        
        # ── Handle tspans ──
        tspans = el.findall(f'{{{SVG_NS}}}tspan')
        
        if not tspans:
            # Simple text: single textbox
            text_content = (el.text or '').strip()
            if not text_content:
                return
            self._create_textbox(slide, x_val, y_val, text_content,
                                 font_size, font_family, font_weight,
                                 fill_color, text_anchor, opacity,
                                 sx, sy, tx, ty)
        else:
            # Multi-tspan: check if tspans have their own x/y/dy
            has_positioned_tspans = any(
                t.get('x') or t.get('y') or t.get('dy')
                for t in tspans
            )
            
            if has_positioned_tspans:
                # Each tspan is independently positioned → separate textboxes
                current_x, current_y = x_val, y_val
                for tspan in tspans:
                    tspan_text = (tspan.text or '').strip()
                    if not tspan_text:
                        continue
                    
                    # tspan can override x, y, dy
                    if tspan.get('x'):
                        current_x = float(tspan.get('x'))
                    dy = tspan.get('dy', '0')
                    if tspan.get('y'):
                        current_y = float(tspan.get('y'))
                    else:
                        # dy can be em (e.g. "1.5em") or px
                        dy_val = dy.strip()
                        if dy_val.endswith('em'):
                            current_y += float(dy_val[:-2]) * font_size
                        elif dy_val.endswith('%'):
                            current_y += float(dy_val[:-1]) / 100 * font_size
                        elif dy_val:
                            current_y += float(dy_val)
                    
                    ts_font_size = float(tspan.get('font-size')) if tspan.get('font-size') else font_size
                    ts_fill = tspan.get('fill', '') or fill_color
                    ts_font_weight = tspan.get('font-weight', '') or font_weight
                    ts_font_family = tspan.get('font-family', '') or font_family
                    ts_anchor = tspan.get('text-anchor', '') or text_anchor
                    ts_opacity = tspan.get('opacity', '') or opacity
                    
                    self._create_textbox(slide, current_x, current_y, tspan_text,
                                         ts_font_size, ts_font_family, ts_font_weight,
                                         ts_fill, ts_anchor, ts_opacity,
                                         sx, sy, tx, ty)
            else:
                # All tspans at same position → single textbox with multiple paragraphs
                all_text = ''.join((t.text or '') for t in tspans).strip()
                if not all_text:
                    return
                self._create_textbox(slide, x_val, y_val, all_text,
                                     font_size, font_family, font_weight,
                                     fill_color, text_anchor, opacity,
                                     sx, sy, tx, ty,
                                     tspan_list=tspans)

    def _create_textbox(self, slide, x_val, y_val, text_content,
                        font_size, font_family, font_weight,
                        fill_color, text_anchor, opacity,
                        sx, sy, tx, ty, tspan_list=None):
        """Create a single PPTX textbox from SVG text parameters."""
        
        # ── Calculate position ──
        x = int(x_val * sx) + tx
        y = int(y_val * sy) + ty
        
        # ── SVG y is baseline position → convert to textbox top ──
        # Approximate: baseline ≈ top + fontSize * 0.8 (for most fonts)
        baseline_offset = int(font_size * 0.8 * PX_TO_EMU * (sx / PX_TO_EMU if sx != PX_TO_EMU else 1))
        # Simpler: just shift up by fontSize * 0.78 in EMU
        baseline_offset = int(font_size * 0.78 * sx)
        y -= baseline_offset
        
        # ── Estimate text width ──
        # Use character count × approximate char width
        cjk_count = self._cjk_count(text_content)
        ascii_count = len(text_content) - cjk_count
        # CJK chars ≈ font_size wide, ASCII chars ≈ font_size * 0.55 wide
        text_width_px = cjk_count * font_size * 1.05 + ascii_count * font_size * 0.6
        text_width_px = max(text_width_px, font_size * 2)  # minimum
        
        # ── Handle text-anchor ──
        if text_anchor == 'middle':
            # SVG x is center point → shift textbox left by half width
            x -= int(text_width_px * sx / 2)
        elif text_anchor == 'end':
            # SVG x is right edge → shift textbox left by full width
            x -= int(text_width_px * sx)
        
        # ── Calculate textbox dimensions ──
        # Width: estimated text width, but capped by slide bounds
        padding = int(8 * PX_TO_EMU)
        w = int(text_width_px * sx) + padding
        # Cap by available space to right edge
        max_w = SLIDE_W - x - padding
        if w > max_w:
            w = max_w
        w = max(w, int(font_size * 2 * sx))  # minimum width
        
        # Height: based on text lines
        line_height = font_size * 1.35
        num_lines = max(1, len(text_content) // max(1, int(text_width_px / (font_size * 0.8))) + 1)
        h = max(int(num_lines * line_height * sx), int(font_size * 1.8 * sx))
        
        # ── Create text box ──
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Zero margins for precise positioning
        from pptx.util import Emu as E
        tf.margin_left = E(0)
        tf.margin_right = E(0)
        tf.margin_top = E(0)
        tf.margin_bottom = E(0)
        
        # ── Populate text ──
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text_content
        run.font.size = Pt(font_size)
        # Extract first font from comma-separated list
        primary_font = font_family.split(',')[0].strip().strip('"\'')
        run.font.name = primary_font
        # Set East Asian font
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            from lxml import etree
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', primary_font)
        
        run.font.bold = font_weight in ('bold', '700', '800', '900')
        fc = parse_color(fill_color)
        if fc:
            run.font.color.rgb = fc
        
        # Alignment
        if text_anchor == 'middle':
            p.alignment = 1  # PP_ALIGN.CENTER
        elif text_anchor == 'end':
            p.alignment = 3  # PP_ALIGN.RIGHT
        
        # ── Apply opacity ──
        if opacity:
            try:
                op_val = float(opacity)
                if op_val < 1.0:
                    alpha = int(op_val * 100000)
                    from lxml import etree
                    solidFill = run._r.get_or_add_rPr().find(qn('a:solidFill'))
                    if solidFill is None:
                        # Need to find the color element in rPr
                        for child in run._r.get_or_add_rPr():
                            if child.tag.endswith('}solidFill') or 'solidFill' in child.tag:
                                solidFill = child
                                break
                    if solidFill is not None:
                        srgb = solidFill.find(qn('a:srgbClr'))
                        if srgb is not None:
                            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                            alpha_el.set('val', str(alpha))
            except (ValueError, TypeError):
                pass
    
    def _add_line(self, slide, el, sx, sy, tx, ty):
        x1 = int(float(el.get('x1', '0')) * sx) + tx
        y1 = int(float(el.get('y1', '0')) * sy) + ty
        x2 = int(float(el.get('x2', '0')) * sx) + tx
        y2 = int(float(el.get('y2', '0')) * sy) + ty
        
        # Use freeform to draw a precise line
        builder = slide.shapes.build_freeform(x1, y1)
        builder.add_line_segments([(x2, y2)])
        
        try:
            shape = builder.convert_to_shape()
            self._apply_stroke(shape, el)
            # Lines shouldn't have fill
            shape.fill.background()
        except Exception:
            # Fallback: use a shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                min(x1, x2), min(y1, y2),
                abs(x2 - x1) or 1, abs(y2 - y1) or 1
            )
            self._apply_stroke(shape, el)
            shape.fill.background()
    
    def _add_path(self, slide, el, sx, sy, tx, ty):
        d = el.get('d', '')
        if not d:
            return
        
        commands = parse_path_d(d)
        if not commands:
            return
        
        xs = [c['x'] for c in commands if c['type'] in ('M', 'L', 'C')]
        ys = [c['y'] for c in commands if c['type'] in ('M', 'L', 'C')]
        if not xs or not ys:
            return
        
        first = commands[0]
        start_x = int(first['x'] * sx) + tx
        start_y = int(first['y'] * sy) + ty
        
        builder = slide.shapes.build_freeform(start_x, start_y)
        
        for cmd in commands[1:]:
            if cmd['type'] == 'L':
                builder.add_line_segments([
                    (int(cmd['x'] * sx) + tx, int(cmd['y'] * sy) + ty)
                ])
            elif cmd['type'] == 'C':
                builder.add_line_segments(
                    [(int(cmd['x'] * sx) + tx, int(cmd['y'] * sy) + ty)],
                    [(int(cmd['cx1'] * sx) + tx, int(cmd['cy1'] * sy) + ty),
                     (int(cmd['cx2'] * sx) + tx, int(cmd['cy2'] * sy) + ty)]
                )
            elif cmd['type'] == 'Z':
                builder.add_line_segments([(start_x, start_y)])
        
        try:
            shape = builder.convert_to_shape()
            self._apply_fill(shape, el)
            self._apply_stroke(shape, el)
        except Exception as e:
            print(f"  ⚠ Path conversion failed: {e}")
    
    def _add_poly(self, slide, el, sx, sy, tx, ty, close=True):
        points_str = el.get('points', '')
        if not points_str:
            return
        
        coords = re.findall(r'[-+]?(?:\d+\.?\d*|\.\d+)', points_str)
        if len(coords) < 4:
            return
        
        pairs = [(float(coords[i]), float(coords[i+1])) for i in range(0, len(coords), 2)]
        
        start_x = int(pairs[0][0] * sx) + tx
        start_y = int(pairs[0][1] * sy) + ty
        
        builder = slide.shapes.build_freeform(start_x, start_y)
        for px_v, py_v in pairs[1:]:
            builder.add_line_segments([(int(px_v * sx) + tx, int(py_v * sy) + ty)])
        if close:
            builder.add_line_segments([(start_x, start_y)])
        
        try:
            shape = builder.convert_to_shape()
            self._apply_fill(shape, el)
            self._apply_stroke(shape, el)
        except Exception as e:
            print(f"  ⚠ Polygon conversion failed: {e}")
    
    def save(self, output_path: str):
        self.prs.save(output_path)
        print(f"  ✅ Saved: {output_path}")


def convert_svg_files(svg_paths: list[str], output_path: str):
    converter = SvgToPptx()
    for i, path in enumerate(svg_paths):
        print(f"  [{i+1}/{len(svg_paths)}] Converting: {path}")
        with open(path, 'r', encoding='utf-8') as f:
            converter.add_svg_slide(f.read())
    converter.save(output_path)


def convert_svg_string(svg_string: str, output_path: str):
    converter = SvgToPptx()
    converter.add_svg_slide(svg_string)
    converter.save(output_path)


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python svg2pptx.py output.pptx input1.svg [input2.svg ...]")
        sys.exit(1)
    convert_svg_files(sys.argv[2:], sys.argv[1])
