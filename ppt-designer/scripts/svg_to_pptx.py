#!/usr/bin/env python3
"""
SVG to PPTX converter.
Usage: python3 svg_to_pptx.py <output.pptx> <slide1.svg> [slide2.svg ...]

Converts SVG slides into a PPTX file, each SVG as one slide (16:9, 1280x720).
Requires: python-pptx (pip install python-pptx)
"""

import sys
import os
import shutil
import subprocess
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

def _run_quiet(cmd):
    try:
        subprocess.run(
            cmd,
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        return True
    except Exception:
        return False


def _render_svg_to_png(temp_svg, svg_content, temp_png):
    # Prefer dedicated SVG renderers when present.
    rsvg = shutil.which("rsvg-convert")
    if rsvg and _run_quiet([rsvg, "-w", "1280", "-h", "720", str(temp_svg), "-o", str(temp_png)]):
        return True

    # macOS ships with `sips`, which can rasterize SVG without extra setup.
    sips = shutil.which("sips")
    if sips and _run_quiet([sips, "-s", "format", "png", str(temp_svg), "--out", str(temp_png)]):
        return True

    try:
        import cairosvg

        cairosvg.svg2png(
            bytestring=svg_content.encode(),
            write_to=str(temp_png),
            output_width=1280,
            output_height=720,
        )
        return True
    except Exception:
        return False


def svg_to_pptx(output_path, svg_files):
    prs = Presentation()
    # Set 16:9 slide size (1280x720 at 96 DPI)
    prs.slide_width = Emu(12192000)   # 10 inches * 914400 EMU/inch
    prs.slide_height = Emu(6858000)   # 7.5 inches * 914400 EMU/inch

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for svg_path in svg_files:
        slide = prs.slides.add_slide(blank_layout)

        # Convert SVG to EMU for sizing
        # 1280px SVG maps to slide width, 720px maps to slide height
        svg_content = Path(svg_path).read_text(encoding='utf-8')

        # Save SVG to temp file for embedding
        temp_svg = Path(output_path).parent / f"_temp_{Path(svg_path).name}"
        temp_svg.write_text(svg_content, encoding='utf-8')

        # Add SVG as image (note: python-pptx doesn't natively support SVG)
        # We'll convert SVG to PNG first using rsvg-convert if available
        temp_png = Path(output_path).parent / f"_temp_{Path(svg_path).stem}.png"

        if _render_svg_to_png(temp_svg, svg_content, temp_png) and temp_png.exists():
            slide.shapes.add_picture(
                str(temp_png),
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height
            )
            temp_png.unlink(missing_ok=True)
        else:
            # Last resort: add text placeholder.
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(8), Inches(1.5)
            )
            tf = txBox.text_frame
            tf.text = f"[SVG Slide: {Path(svg_path).name}]"
            print(
                f"WARNING: Could not render {svg_path}. Install rsvg-convert, use macOS sips, or fix cairosvg/cairo.",
                file=sys.stderr,
            )

        temp_svg.unlink(missing_ok=True)
        print(f"  Added: {svg_path}")

    prs.save(output_path)
    print(f"\nSaved: {output_path} ({len(svg_files)} slides)")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <output.pptx> <slide1.svg> [slide2.svg ...]")
        sys.exit(1)
    svg_to_pptx(sys.argv[1], sys.argv[2:])
